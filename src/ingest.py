from pathlib import Path
from typing import List, Dict, Tuple, Optional
import hashlib

import pdfplumber
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter

from .config import INDEX_DIR, EMBEDDING_MODEL, ALLOWED_EXTS
from .db import upsert_document, insert_chunks
from .topics import process_and_store_topics


def _hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]


def extract_text_from_pdf(path: Path) -> List[Tuple[str, Optional[int]]]:
    texts: List[Tuple[str, Optional[int]]] = []
    try:
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                content = page.extract_text() or ""
                if content.strip():
                    texts.append((content, i))
    except Exception:
        # Fallback: PyPDF2
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, start=1):
            content = page.extract_text() or ""
            if content.strip():
                texts.append((content, i))
    return texts


def extract_text_from_docx(path: Path) -> List[Tuple[str, Optional[int]]]:
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    return [("\n".join(paragraphs), None)]


def extract_text_from_pptx(path: Path) -> List[Tuple[str, Optional[int]]]:
    prs = Presentation(str(path))
    slides_text: List[Tuple[str, Optional[int]]] = []
    for i, slide in enumerate(prs.slides, start=1):
        buf: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                buf.append(shape.text)
        content = "\n".join(buf)
        if content.strip():
            slides_text.append((content, i))
    return slides_text


def _split_chunks(texts: List[Tuple[str, Optional[int]]]) -> List[Tuple[str, Optional[int]]]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", ", ", " "]
    )
    chunks: List[Tuple[str, Optional[int]]] = []
    for content, page in texts:
        for ch in splitter.split_text(content):
            chunks.append((ch, page))
    return chunks


def ingest_files(files: List[Path], subject: Optional[str] = None, doc_type: Optional[str] = None, year: Optional[int] = None, tags: Optional[str] = None) -> Dict:
    if not files:
        return {"ok": False, "error": "No files selected"}

    # Validate extensions
    for f in files:
        if f.suffix.lower() not in ALLOWED_EXTS:
            return {"ok": False, "error": f"Unsupported file type: {f.name}"}

    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    texts_all: List[str] = []
    metas_all: List[Dict] = []

    for f in files:
        path = Path(f)
        ext = path.suffix.lower()
        name = path.name

        meta = {
            "path": str(path),
            "name": name,
            "ext": ext,
            "subject": subject,
            "doc_type": doc_type,
            "year": year,
            "tags": tags,
        }
        doc_id = upsert_document(meta)

        if ext == ".pdf":
            texts = extract_text_from_pdf(path)
        elif ext == ".docx":
            texts = extract_text_from_docx(path)
        elif ext == ".pptx":
            texts = extract_text_from_pptx(path)
        else:
            texts = []

        chunks = _split_chunks(texts)
        # Insert chunks per document
        insert_chunks(doc_id=doc_id, chunks=[(_hash(content), content, page) for content, page in chunks])
        # Gather for index and topic extraction
        doc_texts = []
        for content, page in chunks:
            texts_all.append(content)
            doc_texts.append(content)
            metas_all.append({
                "doc_id": doc_id,
                "name": name,
                "subject": subject,
                "doc_type": doc_type,
                "year": year,
                "page": page,
            })
        # Topic extraction for this document
        try:
            process_and_store_topics(doc_id, doc_texts, top_n=20)
        except Exception:
            # If spaCy model missing, skip gracefully
            pass

    # Build or update FAISS index
    index_path = str(INDEX_DIR)
    if (INDEX_DIR / "index.faiss").exists():
        vs = FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
        vs.add_texts(texts_all, metas_all)
        vs.save_local(index_path)
    else:
        vs = FAISS.from_texts(texts_all, embeddings, metadatas=metas_all)
        vs.save_local(index_path)

    return {"ok": True, "indexed": len(texts_all), "documents": len(files)}
